#!/usr/bin/env python3.10
# -*- coding: utf-8 -*-


import os
import tempfile
import argparse
import subprocess
import shutil
from subprocess import call
from PIL import Image

from pdf2image import convert_from_path
from gtts import gTTS
from pptx import Presentation
from TTS.utils.manage import ModelManager
from TTS.utils.synthesizer import Synthesizer
import site


__author__ = ['chaonan99','Mathieu']

def ppt_presenter(pptx_path, language, use_gtts, output_path):
    if use_gtts:
        print("Using gtts\n")
    else:
        print("Using TTS\n")
        location = site.getsitepackages()[0]
        path = location + "/TTS/.models.json"
        model_manager = ModelManager(path)
        if language == "fr":
            model_path, config_path, model_item = model_manager.download_model("tts_models/fr/css10/vits")
            syn = Synthesizer(
                tts_checkpoint=model_path,
                tts_config_path=config_path,

            )
        else:
            model_path, config_path, model_item = model_manager.download_model("tts_models/en/ljspeech/tacotron2-DDC")
            voc_path, voc_config_path, _ = model_manager.download_model(model_item["default_vocoder"])

            syn = Synthesizer(
                tts_checkpoint=model_path,
                tts_config_path=config_path,
                vocoder_checkpoint=voc_path,
                vocoder_config=voc_config_path
            )


    with tempfile.TemporaryDirectory(dir=os.getcwd()) as temp_path:
        filename = os.path.splitext(os.path.basename(pptx_path))[0]
        pdf_path = os.path.join(temp_path, filename + '.pdf')
        print('Converting pptx to pdf: ' + pdf_path + '\n')
        call(['soffice', '--headless', '--convert-to', 'pdf', '--outdir' , temp_path, '--convert-images-to', 'jpg', pptx_path])
        print('pdf created\n')
        # temp_path = "./temp"
        images_from_path = convert_from_path(pdf_path)
        prs = Presentation(pptx_path)
        assert len(images_from_path) == len(prs.slides)
        count = len(prs.slides)
        for i, (slide, image) in enumerate(zip(prs.slides, images_from_path)):
            print('Processing {}/{}\n'.format(i + 1, count))
            if slide.has_notes_slide:
                image_path = os.path.join(temp_path, 'frame_{}.jpg'.format(i))
               
                image.save(image_path)
                make_jpeg_even(image_path)
                notes = slide.notes_slide.notes_text_frame.text
                if use_gtts:
                    audio_path = os.path.join(temp_path, 'frame_{}.mp3'.format(i))
                    tts = gTTS(text=notes, lang=language)
                    tts.save(audio_path)
                else:
                    audio_path = os.path.join(temp_path, 'frame_{}.wav'.format(i))
                    outputs = syn.tts(notes)
                    syn.save_wav(outputs, audio_path)

                ffmpeg_call(image_path, audio_path, temp_path, i)
        print('Concatenating videos...\n')
        video_list = [os.path.join(temp_path, 'frame_{}.ts'.format(i)) \
                        for i in range(len(images_from_path))]
        video_list_str = 'concat:' + '|'.join(video_list)
        ffmpeg_concat(video_list_str, output_path)
        print('Generation completed\n')

def make_jpeg_even(image_path):
    img = Image.open(image_path)

    # Crop image to have even dimensions
    even_width = img.size[0] - (img.size[0] % 2)
    even_height = img.size[1] - (img.size[1] % 2)
    img = img.crop((0, 0, even_width, even_height))

    # Save cropped image
    img.save(image_path)


def ffmpeg_call(image_path, audio_path, temp_path, i):
    out_path_mp4 = os.path.join(temp_path, 'frame_{}.mp4'.format(i))
    out_path_ts = os.path.join(temp_path, 'frame_{}.ts'.format(i))
    print('Checking audio duration...\n')
    duration = subprocess.check_output([
        'ffprobe', 
        '-i', audio_path, 
        '-show_entries', 
        'format=duration', 
        '-v', 'quiet', 
        '-of', 'csv=p=0'])
    duration = float(duration)
    print("Audio duration: " + str(duration) + "\n")
    call(['ffmpeg', '-loop', '1', '-y', '-i', image_path, '-i', audio_path,
          '-c:v', 'libx264', '-tune', 'stillimage', '-c:a', 'aac',
          '-b:a', '192k', '-pix_fmt', 'yuv420p', '-t', str(duration) , out_path_mp4])
    print('Converting to ts...\n')
    call(['ffmpeg', '-y', '-i', out_path_mp4, '-c', 'copy',
          '-bsf:v', 'h264_mp4toannexb', '-f', 'mpegts', out_path_ts])


def ffmpeg_concat(video_list_str, out_path):
    call(['ffmpeg', '-y', '-f', 'mpegts', '-i', '{}'.format(video_list_str),
          '-c', 'copy', '-bsf:a', 'aac_adtstoasc', out_path])

def is_executable_in_path(executable_name):
    return shutil.which(executable_name) is not None

def main():
    parser = argparse.ArgumentParser(description='PPTX 2 mp4 help.')
    parser.add_argument('-i', '--input', help='input pptx path')
    parser.add_argument('-l', '--lang', help='text language (fr, en, etc...)')
    parser.add_argument('-o', '--output', help='output path')
    parser.add_argument('-t', '--tts', help='Text to speech engine to use (gtts or tts). Default is tts (CoquiTTS)', required=False, default='tts')
    args = parser.parse_args()

    error = False
    if not is_executable_in_path('soffice'):
        print('Please, install LibreOffice and ensure "soffice" is in your path.\n')
        error = True
    if not is_executable_in_path('ffmpeg'):
        print('Please, install ffmpeg and ensure "ffmpeg" is in your path.\n')
        error = True
    if error:
        print('Error: missing software. Please install required software and try again.\n')
        exit(1)

    ppt_presenter(args.input, args.lang, args.tts == 'gtts', args.output)


if __name__ == '__main__':
    main()